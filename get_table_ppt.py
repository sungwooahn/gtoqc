import cirq
import numpy as np
from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # Class for text frame positioning
from pptx.util import Pt 

from cirq.contrib.qasm_import import circuit_from_qasm
from parse import compile
from optparse import OptionParser
from tabulate import tabulate

#### usage format
######## python3 get_table_ppt.py -b bv -q 5

parser = OptionParser()
parser.add_option("-g", "--use_gpu", dest="use_gpu", default=True, help="True for GPU usage")
parser.add_option("-b", "--bench_name", dest="bench_name", help="set benchmark name")
parser.add_option("-q", "--num_qubits", dest="num_qubits", help="number of qubits in benchmark")

(options, args) = parser.parse_args()

### file size
# bv: 10 32 2
# adder 10 31 3
# grover: 11 31 2
# pe: 10 30 2
# qft: 10 32 2

USE_GPU=options.use_gpu
if USE_GPU=="True":
    USE_GPU=True
else:
    USE_GPU=False

bench_name=options.bench_name
bench_dir="/home/ahnsungwoo/blinq/blinq_bench"
bench_type=bench_name

num_qubits=options.num_qubits
bench_qnum=num_qubits

bench_postfix=".qasm"
file_path=bench_dir+"/"+bench_type+"/"+bench_type+"_"+bench_qnum+bench_postfix
bench = open(file_path,mode='r')
qasm_input = bench.read()
bench.close()


class sw_Gate(cirq.Gate):
    def __init__(self, sw_qubits, sw_unitary, sw_str):
        self._qubits=sw_qubits
        self._unitary=sw_unitary
        self._str=sw_str
    def _num_qubits_(self) -> int:
        return self._qubits

    def _unitary_(self):
        return self._unitary

    def __str__(self):
        return self._str

@cirq.transformer
class SerialMergeFunc:
    def __init__(self, degree):
        self._degree = degree
        
    def find_second_candidate(self, circuit, first_op, first_moment, second_list):
        second_candidate_moment = circuit.next_moment_operating_on(first_op.qubits, first_moment+1)
        if second_candidate_moment == None:
            return second_list
        else:
            second_candidate_op = circuit.operation_at(first_op.qubits[0], second_candidate_moment)
            if second_candidate_op == None:
                return second_list
            elif second_candidate_op.qubits==first_op.qubits:
                second_list.append(second_candidate_op)
                second_list = self.find_second_candidate(circuit, second_candidate_op, second_candidate_moment, second_list)
                return second_list
            
            else:
                return second_list
                
    def check_merged(self, circuit, target_op, target_moment):
        merge_status = False
        previous_moment_list = []
        previous_op_list = []
        checking_length = len(target_op.qubits)
        for index in range(0, checking_length):
            checking_qubit = [target_op.qubits[index]]
            prev_moment_index = circuit.prev_moment_operating_on(checking_qubit, target_moment)
            if prev_moment_index == None:
                return False
            else:
                previous_moment_list.append(prev_moment_index)

        check_prev_moment_value = previous_moment_list[0]
        for index in range(0, len(previous_moment_list)):
            if previous_moment_list[index] != check_prev_moment_value:
                return False
            
            previous_op = circuit.operation_at(target_op.qubits[index], previous_moment_list[index])
            previous_op_list.append(previous_op)
            
        for index in range(0, len(previous_op_list)):
            if previous_op_list[index].qubits != target_op.qubits:
                return False
                
        merge_status = True    
        return True
            

            

    def __call__(self, circuit, *, context=None):
        merged_circuit = cirq.Circuit()
        all_op = circuit.findall_operations(lambda op:op)
        for i, op in all_op:
            first_candidiate_op = op
            if len(first_candidiate_op.qubits)==self._degree:
                if self.check_merged(circuit, first_candidiate_op, i) == True:
                    continue
                second_list = []
                second_list = self.find_second_candidate(circuit, first_candidiate_op, i, second_list)
                if len(second_list) !=0:
                    new_unitary = cirq.unitary(first_candidiate_op)
                    # new_unitary = cirq.kraus(first_candidiate_op)
                    for op in second_list:
                        new_unitary = np.matmul(cirq.unitary(op), new_unitary)
                        # new_unitary = np.matmul(cirq.kraus(op), new_unitary)
                    new_gate = sw_Gate(self._degree, new_unitary, 'sm')
                    new_qubit = [qubits for qubits in first_candidiate_op.qubits]
                    merged_circuit.append(new_gate(*new_qubit))
                    
                else:
                    merged_circuit.append(first_candidiate_op)
                
            else:
                merged_circuit.append(first_candidiate_op)
        
        # print(merged_circuit)
        return merged_circuit

@cirq.transformer
class SerialMergeFunc_moment:
    def __init__(self, degree):
        self._degree = degree
        
    def find_second_candidate(self, circuit, first_op, first_moment, second_list):
        second_candidate_moment = circuit.next_moment_operating_on(first_op.qubits, first_moment+1)
        if second_candidate_moment == None:
            return second_list
        else:
            second_candidate_op = circuit.operation_at(first_op.qubits[0], second_candidate_moment)
            if second_candidate_op == None:
                return second_list
            elif second_candidate_op.qubits==first_op.qubits:
                second_list.append(second_candidate_op)
                second_list = self.find_second_candidate(circuit, second_candidate_op, second_candidate_moment, second_list)
                return second_list
            
            else:
                return second_list
                
    def check_merged(self, circuit, target_op, target_moment):
        merge_status = False
        previous_moment_list = []
        previous_op_list = []
        checking_length = len(target_op.qubits)
        for index in range(0, checking_length):
            checking_qubit = [target_op.qubits[index]]
            prev_moment_index = circuit.prev_moment_operating_on(checking_qubit, target_moment)
            if prev_moment_index == None:
                return False
            else:
                previous_moment_list.append(prev_moment_index)

        check_prev_moment_value = previous_moment_list[0]
        for index in range(0, len(previous_moment_list)):
            if previous_moment_list[index] != check_prev_moment_value:
                return False
            
            previous_op = circuit.operation_at(target_op.qubits[index], previous_moment_list[index])
            previous_op_list.append(previous_op)
            
        for index in range(0, len(previous_op_list)):
            if previous_op_list[index].qubits != target_op.qubits:
                return False
                
        merge_status = True    
        return True
            

    def __call__(self, circuit, *, context=None):
        merged_circuit = cirq.Circuit()
        # all_op = circuit.findall_operations(lambda op:op)
        
        moment_index = -1
        for moment in circuit:
            moment_index += 1
            append_moment=[]
            for op in moment:
                first_candidiate_op = op
                if len(first_candidiate_op.qubits)==self._degree:
                    if self.check_merged(circuit, first_candidiate_op, moment_index) == True:
                        continue
                    second_list = []
                    second_list = self.find_second_candidate(circuit, first_candidiate_op, moment_index, second_list)
                    if len(second_list) !=0:
                        new_unitary = cirq.unitary(first_candidiate_op)
                        # new_unitary = cirq.kraus(first_candidiate_op)
                        for op in second_list:
                            new_unitary = np.matmul(cirq.unitary(op), new_unitary)
                            # new_unitary = np.matmul(cirq.kraus(op), new_unitary)
                        new_gate = sw_Gate(self._degree, new_unitary, 'sm')
                        new_qubit = [qubits for qubits in first_candidiate_op.qubits]
                        append_moment.append(new_gate(*new_qubit))
                    else:
                        append_moment.append(first_candidiate_op)
                        
                else:
                    append_moment.append(first_candidiate_op)
                    # merged_circuit.append(first_candidiate_op)
                    
            merged_circuit.append(append_moment, strategy=cirq.InsertStrategy.NEW_THEN_INLINE)
            # print(merged_circuit)
                
        
        # for i, op in all_op:
        #     first_candidiate_op = op
        #     if len(first_candidiate_op.qubits)==self._degree:
        #         if self.check_merged(circuit, first_candidiate_op, i) == True:
        #             continue
        #         second_list = []
        #         second_list = self.find_second_candidate(circuit, first_candidiate_op, i, second_list)
        #         if len(second_list) !=0:
        #             new_unitary = cirq.unitary(first_candidiate_op)
        #             # new_unitary = cirq.kraus(first_candidiate_op)
        #             for op in second_list:
        #                 new_unitary = np.matmul(cirq.unitary(op), new_unitary)
        #                 # new_unitary = np.matmul(cirq.kraus(op), new_unitary)
        #             new_gate = sw_Gate(self._degree, new_unitary, 'sm')
        #             new_qubit = [qubits for qubits in first_candidiate_op.qubits]
        #             merged_circuit.append(new_gate(*new_qubit))
                    
        #         else:
        #             merged_circuit.append(first_candidiate_op)
                
        #     else:
        #         merged_circuit.append(first_candidiate_op)
        
        # print(merged_circuit)
        return merged_circuit

@cirq.transformer
class SerialMergeFunc_collapse:
    def __init__(self, degree):
        self._degree = degree
        
    def construct_collapsed_circuit_left(self, circuit, first_moment_index, first_op, second_candidate_moment_index, second_candidate_op):  
        collapsed_circuit = cirq.Circuit()
        
        for moment in circuit[:second_candidate_moment_index]:
            collapsed_circuit.append(moment)
            
        target_moment = circuit[second_candidate_moment_index]
        replaced_moment = target_moment.without_operations_touching(second_candidate_op.qubits)
        collapsed_circuit.append(replaced_moment)
            
        for moment in circuit[second_candidate_moment_index+1:]:
            collapsed_circuit.append(moment)
        
        # print(first_op)
        # print(second_candidate_op)
        # print(collapsed_circuit)
        return collapsed_circuit
    
    def check_collapse_possible_left(self, circuit, first_moment_index, first_op, second_candidate_moment_index, second_candidate_op):
        for second_index in range(len(second_candidate_op.qubits)):
            check_op_index = circuit.next_moment_operating_on([second_candidate_op.qubits[second_index]], second_candidate_moment_index+1)
            if check_op_index != first_moment_index:
                return False
            if first_op != circuit.operation_at(second_candidate_op.qubits[second_index], check_op_index):
                return False
            
        return True
    
    def collapse_left(self, circuit, first_moment_index, first_op):
        
        for index in range(len(first_op.qubits)):
            second_candidate_moment_index = circuit.prev_moment_operating_on([first_op.qubits[index]], first_moment_index)
            if second_candidate_moment_index != None:
                second_candidate_op = circuit.operation_at(first_op.qubits[index], second_candidate_moment_index)
                collapse_possible = self.check_collapse_possible_left(circuit, first_moment_index, first_op, second_candidate_moment_index, second_candidate_op)
                
                if collapse_possible:
                    collapsed_circuit = self.construct_collapsed_circuit_left(circuit, first_moment_index, first_op, second_candidate_moment_index, second_candidate_op)
                    return collapsed_circuit, True
        
        return circuit, False


    def construct_collapsed_circuit_right(self, circuit, first_moment_index, first_op, second_candidate_moment_index, second_candidate_op):  
        collapsed_circuit = cirq.Circuit()
        
        for moment in circuit[:second_candidate_moment_index]:
            collapsed_circuit.append(moment)
            
        target_moment = circuit[second_candidate_moment_index]
        replaced_moment = target_moment.without_operations_touching(second_candidate_op.qubits)
        collapsed_circuit.append(replaced_moment)
            
            
        if len(circuit) > second_candidate_moment_index+1:
            for moment in circuit[second_candidate_moment_index+1:]:
                collapsed_circuit.append(moment)
        
        return collapsed_circuit
    def check_collapse_possible_right(self, circuit, first_moment_index, first_op, second_candidate_moment_index, second_candidate_op):
        
        for second_index in range(len(second_candidate_op.qubits)):
            check_op_index = circuit.prev_moment_operating_on([second_candidate_op.qubits[second_index]], second_candidate_moment_index)
            if check_op_index != first_moment_index:
                return False
            if first_op != circuit.operation_at(second_candidate_op.qubits[second_index], check_op_index):
                return False
            
        return True
    def collapse_right(self, circuit, first_moment_index, first_op):
        
        for index in range(len(first_op.qubits)):
            second_candidate_moment_index = circuit.next_moment_operating_on([first_op.qubits[index]], first_moment_index+1)
            if second_candidate_moment_index != None:
                second_candidate_op = circuit.operation_at(first_op.qubits[index], second_candidate_moment_index)
                collapse_possible = self.check_collapse_possible_right(circuit, first_moment_index, first_op, second_candidate_moment_index, second_candidate_op)
                
                if collapse_possible:
                    collapsed_circuit = self.construct_collapsed_circuit_right(circuit, first_moment_index, first_op, second_candidate_moment_index, second_candidate_op)
                    return collapsed_circuit, True
        
        return circuit, False
        
            
    def merge_collapse(self, circuit):
        
        changed = False
        
        all_op = circuit.findall_operations(lambda op:op)
        for first_moment_index, first_op in all_op:
            circuit, changed = self.collapse_left(circuit, first_moment_index, first_op)
            if changed:
                # print(circuit)
                return circuit, True
            
            circuit, changed = self.collapse_right(circuit, first_moment_index, first_op)
            if changed:
                # print(circuit)
                return circuit, True
        
        return circuit, False
        
    def __call__(self, circuit, *, context=None):
        collapsed_circuit=cirq.Circuit()
        changed = True
        while changed == True:
            circuit, changed = self.merge_collapse(circuit)
        
        for moment in circuit:
            if len(moment)!=0:
                collapsed_circuit.append(moment)
            
        # merged_circuit = circuit
        return collapsed_circuit

@cirq.transformer
class ParallelMergeFunc_full:
    def __init__(self, degree):
        self._degree = degree
            

    
    def get_coefficient(self, op1, op2, op1_qubit_index, op2_qubit_index, c_row, c_col):
        
        # required qubit index list
        # op_index_original
        # op_index_shrinked_inter
        # op_index_shrinked_intra
        
        op1_index_original = op1_qubit_index
        op2_index_original = op2_qubit_index
        
        # if (op2_qubit_index==[2, 3, 4]):
        #     print(op2)
        
        op1_index_shrinked_inter = []
        op2_index_shrinked_inter = []
        total_index = sorted(op1_index_original + op2_index_original)
        op1_len = len(op1_qubit_index)
        op2_len = len(op2_qubit_index)
        
        for index in range(len(total_index)):
            if total_index[index] in op1_index_original:
                op1_index_shrinked_inter.append(index)
            else:
                op2_index_shrinked_inter.append(index)
        
        op1_index_shrinked_intra = []
        op2_index_shrinked_intra = []
        
        for index in range(op1_len):
            op1_index_shrinked_intra.append(index)
        for index in range(op2_len):
            op2_index_shrinked_intra.append(index)
                
        # if (op2_qubit_index==[2, 3, 4]):
        #     print(op2_index_shrinked_inter)        
        #     print(op2_index_shrinked_intra) 
        #     print()       
                
        # c_row represents the actual final state itself
        # the result of this "get_coefficient" represents the value from "c_col" state to "c_row"
        # transition represents c_col to c_row
        c_col_bin = [int(x) for x in bin(c_col)[2:]]
        c_row_bin = [int(x) for x in bin(c_row)[2:]]
        
        while (len(c_col_bin) != len(total_index)):
            c_col_bin.insert(0, 0)
        while (len(c_row_bin) != len(total_index)):
            c_row_bin.insert(0, 0)
        
        op1_row_bin = []
        op1_col_bin = []
        
        for index in op1_index_shrinked_inter:
            op1_row_bin.append(c_row_bin[index])
            op1_col_bin.append(c_col_bin[index])
            
        op2_row_bin = []
        op2_col_bin = []
        
        
        for index in op2_index_shrinked_inter:
            op2_row_bin.append(c_row_bin[index])
            op2_col_bin.append(c_col_bin[index])
        
        
        op1_row_coefficient = 0
        op1_col_coefficient = 0
        
        for index in range(op1_len):
            op1_row_coefficient += op1_row_bin[index] << op1_index_shrinked_intra[op1_len-index-1]
            op1_col_coefficient += op1_col_bin[index] << op1_index_shrinked_intra[op1_len-index-1]
            
        op2_row_coefficient = 0
        op2_col_coefficient = 0
        
        for index in range(op2_len):
            op2_row_coefficient += op2_row_bin[index] << op2_index_shrinked_intra[op2_len-index-1] 
            op2_col_coefficient += op2_col_bin[index] << op2_index_shrinked_intra[op2_len-index-1] 
        
        op1_coefficient = cirq.unitary(op1)[op1_row_coefficient][op1_col_coefficient]
        op2_coefficient = cirq.unitary(op2)[op2_row_coefficient][op2_col_coefficient]
        
        
        final_coefficient = op1_coefficient * op2_coefficient
        
        # if c_row==5 and c_col==7:
        #     if (op2_qubit_index==[2, 3, 4]):
        #         # print(final_coefficient)
        #         print("index: ", c_row, c_col)
        #         print(op1_row_coefficient)
        #         print(op1_col_coefficient)
        #         print(op2_row_coefficient)
        #         print(op2_col_coefficient)
        #         print(op2_coefficient)
        #         # print(cirq.unitary(op2))
                
                # print()
        
        return final_coefficient
    
    def make_clustered_op(self, op1, op2):
        
        
        op1_qubits = [qubits for qubits in op1.qubits]
        op2_qubits = [qubits for qubits in op2.qubits]
        
        new_op_qubit = sorted(op1_qubits + op2_qubits)
        new_op_unitary_dim = pow(2, len(new_op_qubit))
        
        parse_qubit_index = compile("q_{}")
        # parse_qubit_index = compile("q{}")
        
        op1_qubit_index=[]
        op2_qubit_index=[]
        
        for index in op1_qubits:
            op1_qubit_index.append(int((parse_qubit_index.parse(str(index)))[0]))
        for index in op2_qubits:
            op2_qubit_index.append(int((parse_qubit_index.parse(str(index)))[0]))
        
        # if (op2_qubit_index==[2, 4, 3]):
        #     print(cirq.unitary(op2))
        #     print("euraka")
        # op1_qubit_index=sorted(op1_qubit_index)
        # op2_qubit_index=sorted(op2_qubit_index)
        
        # if (op2_qubit_index==[2, 3, 4]):
        #     print("sorted")
        
        new_op_unitary_row = []
        new_op_unitary = np.empty((0,new_op_unitary_dim), float)   
        
        for c_row in range(new_op_unitary_dim):
            new_op_unitary_row = []
            for c_col in range(new_op_unitary_dim):
                
                final_coefficient = self.get_coefficient(op1, op2, op1_qubit_index, op2_qubit_index, c_row, c_col)
                
                new_op_unitary_row.append(final_coefficient)
            new_op_unitary = np.append(new_op_unitary, np.array([new_op_unitary_row]), axis=0)

        new_gate = sw_Gate(self._degree, new_op_unitary, 'pm')
        new_op = new_gate(*new_op_qubit)
        return new_op
    
        

    def __call__(self, circuit, context=None):
        merged_circuit = cirq.Circuit()
        qubit_tuple = sorted(circuit.all_qubits())
        for moment in circuit:
            new_moment = []
            merged_op_list = []
            non_affecting_op_list=[]
            moment_op_list = []
            for qubit in qubit_tuple:
                for op in moment:
                    if qubit in op.qubits and op not in moment_op_list:
                        moment_op_list.append(op)
            cluster_candidate_op_list = moment_op_list
            
            if(len(cluster_candidate_op_list)>=2):
                for op1_index in range(len(cluster_candidate_op_list)-1):
                    op1_candidate = cluster_candidate_op_list[op1_index]
                    
                    # if op1_candidate in non_affecting_op_list:
                    #     continue
                    
                    for op2_index in range(op1_index+1, len(cluster_candidate_op_list)):
                        op2_candidate = cluster_candidate_op_list[op2_index]
                        
                        # if op2_candidate in non_affecting_op_list:
                        #     continue
                        
                        if op1_candidate in merged_op_list:
                            if (op1_index == len(cluster_candidate_op_list)-2 and op2_index == len(cluster_candidate_op_list)-1):
                                new_moment.append(op2_candidate)
                            continue
                        if op2_candidate in merged_op_list:
                            continue
                        if len(op1_candidate.qubits) + len(op2_candidate.qubits) != self._degree:
                            if (op1_index == len(cluster_candidate_op_list)-2 and op2_index == len(cluster_candidate_op_list)-1):
                                new_moment.append(op2_candidate)
                            continue
                        
                        
                        
                        
                        op1_qubits = [qubits for qubits in op1_candidate.qubits]
                        op2_qubits = [qubits for qubits in op2_candidate.qubits]

                        parse_qubit_index = compile("q_{}")
                        # parse_qubit_index = compile("q{}")

                        op1_qubit_index=[]
                        op2_qubit_index=[]

                        for index in op1_qubits:
                            op1_qubit_index.append(int((parse_qubit_index.parse(str(index)))[0]))
                        for index in op2_qubits:
                            op2_qubit_index.append(int((parse_qubit_index.parse(str(index)))[0]))
                            
                        if op1_qubit_index!=sorted(op1_qubit_index):
                            if op1_candidate in non_affecting_op_list:
                                pass
                            else:
                                new_moment.append(op1_candidate)
                                non_affecting_op_list.append(op1_candidate)
                                
                            if op2_qubit_index!=sorted(op2_qubit_index):
                                if op2_candidate in non_affecting_op_list:
                                    pass
                                else:
                                    new_moment.append(op2_candidate)
                                    non_affecting_op_list.append(op2_candidate)
                            elif op2_index == len(cluster_candidate_op_list)-1:
                                new_moment.append(op2_candidate)
                            else:
                                pass
                            continue
                            
                        if op2_qubit_index!=sorted(op2_qubit_index):
                            if op2_candidate in non_affecting_op_list:
                                pass
                            else:
                                new_moment.append(op2_candidate)
                                non_affecting_op_list.append(op2_candidate)
                            continue
                        
                        
                        
                        
                        merged_op_list.append(op1_candidate)
                        merged_op_list.append(op2_candidate)
                        
                        # print(op1_candidate)
                        # print(op2_candidate)
                        # print(op2_candidate.qubits)
                        # print(cirq.unitary(op2_candidate))
                        new_op=self.make_clustered_op(op1_candidate, op2_candidate)
                    
                        new_moment.append(new_op)
                    
                    if op1_candidate in non_affecting_op_list:
                        continue
                    
                    if op1_candidate not in merged_op_list:
                        new_moment.append(op1_candidate)
                        
                        if op1_index == len(cluster_candidate_op_list)-1:
                            new_moment.append(op2_candidate)
                            
            else:
                new_moment = moment
            merged_circuit.append(new_moment, strategy=cirq.InsertStrategy.NEW_THEN_INLINE)
            # print(merged_circuit)
        return merged_circuit
   
@cirq.transformer
class InsertFrontAndBack_without_expand:
    def __init__(self, degree):
        self._degree = degree
            

    def target_insertion(self, circuit, moment_index):
        current_moment=circuit[moment_index]
        
        current_moment_op_qubit = []
        current_moment_noop_qubit = []
        insert_position = []
    
        all_position = list(circuit.all_qubits())
        
        for qubit in all_position:
            if current_moment.operates_on([qubit]):
                current_moment_op_qubit.append(qubit)
            else:
                current_moment_noop_qubit.append(qubit)
        
        if(moment_index==0):
            if len(circuit) == 1:
                pass
            else:
                next_moment = circuit[moment_index+1]
                for qubit1 in current_moment_noop_qubit:
                    if next_moment.operates_on([qubit1]):
                        insert_position.append(qubit1)

        elif(moment_index==len(circuit)-1):
            prev_moment = circuit[moment_index-1]
            for qubit2 in current_moment_noop_qubit:
                if prev_moment.operates_on([qubit2]):
                    insert_position.append(qubit2)

        else:    
            prev_moment = circuit[moment_index-1]
            next_moment = circuit[moment_index+1]
            for qubit3 in current_moment_noop_qubit:
                if prev_moment.operates_on([qubit3]):
                    insert_position.append(qubit3)

            for qubit4 in current_moment_noop_qubit:
                if next_moment.operates_on([qubit4]):
                    insert_position.append(qubit4)

        insert_position_shrink=[]
        for index in insert_position:
            if index not in insert_position_shrink:
                insert_position_shrink.append(index)

        identity_op_unitary = np.array([[0,1], [1, 0]])   
        inserted_moment = current_moment
        
        # print("start")
        # print(insert_position_shrink)
        for index in insert_position_shrink:
            index_to_list = []
            index_to_list.append(index)
            
            identity_gate = sw_Gate(1, identity_op_unitary, 'iden')
            identity_op = identity_gate(*index_to_list)
            inserted_moment=inserted_moment.with_operation(identity_op)
        return inserted_moment


    def __call__(self, circuit, *, context=None):
        modified_circuit = cirq.Circuit()
        
        moment_index = -1
        for moment in circuit:
            moment_index += 1
            inserted_moment = self.target_insertion(circuit, moment_index)
            modified_circuit.append(inserted_moment)
                     
        return modified_circuit    

### inserts identity by watching odd and even index moments dual in only.
@cirq.transformer
class Insert_dual:
    def __init__(self, degree):
        self._degree = degree
            
    def target_insertion_odd(self, circuit, moment_index):
        current_moment=circuit[moment_index]
        
        current_moment_op_qubit = []
        current_moment_noop_qubit = []
        insert_position = []
    
        all_position = list(circuit.all_qubits())
        
        for qubit in all_position:
            if current_moment.operates_on([qubit]):
                current_moment_op_qubit.append(qubit)
            else:
                current_moment_noop_qubit.append(qubit)
        

        prev_moment = circuit[moment_index-1]
        for qubit3 in current_moment_noop_qubit:
            if prev_moment.operates_on([qubit3]):
                insert_position.append(qubit3)


        identity_op_unitary = np.array([[0,1], [1, 0]])   
        inserted_moment = current_moment
        
        if len(current_moment_op_qubit) + len(insert_position)> self._degree:
            pass
        else:
            for index in insert_position:
                index_to_list = []
                index_to_list.append(index)

                identity_gate = sw_Gate(1, identity_op_unitary, 'I')
                identity_op = identity_gate(*index_to_list)
                inserted_moment=inserted_moment.with_operation(identity_op)
                
        return inserted_moment


    def target_insertion_even(self, circuit, moment_index):
        current_moment=circuit[moment_index]
        
        current_moment_op_qubit = []
        current_moment_noop_qubit = []
        insert_position = []
    
        all_position = list(circuit.all_qubits())
        
        for qubit in all_position:
            if current_moment.operates_on([qubit]):
                current_moment_op_qubit.append(qubit)
            else:
                current_moment_noop_qubit.append(qubit)
        
        next_moment = circuit[moment_index+1]
        for qubit4 in current_moment_noop_qubit:
            if next_moment.operates_on([qubit4]):
                insert_position.append(qubit4)


        identity_op_unitary = np.array([[1,0], [0, 1]])   
        inserted_moment = current_moment
        
        if len(current_moment_op_qubit) + len(insert_position)> self._degree:
            pass
        else:
            for index in insert_position:
                index_to_list = []
                index_to_list.append(index)

                identity_gate = sw_Gate(1, identity_op_unitary, 'I')
                identity_op = identity_gate(*index_to_list)
                inserted_moment=inserted_moment.with_operation(identity_op)
                
        return inserted_moment


    def __call__(self, circuit, *, context=None):
        modified_circuit = cirq.Circuit()
        
        end_moment_index = len(circuit)
        if len(circuit)%2==1:
            end_moment_index=len(circuit)-1
            
        moment_index = -1
        for moment in circuit[:end_moment_index]:
            moment_index += 1
            if moment_index%2==0:
                # print(moment_index)
                inserted_moment = self.target_insertion_even(circuit, moment_index)
            else:
                inserted_moment = self.target_insertion_odd(circuit, moment_index)
            modified_circuit.append(inserted_moment)
                     
        if len(circuit)%2==1:
            modified_circuit.append(circuit[-1])          
                     
        return modified_circuit    

def conduct_clustering(circuit, max_cluster_degree):
    serial_func = SerialMergeFunc(1)
    clustered_circuit = serial_func(circuit)
    
    for degree in range(2,max_cluster_degree+1):
        parallel_func=ParallelMergeFunc_full(degree)
        serial_func=SerialMergeFunc(degree)
        
        clustered_circuit = parallel_func(clustered_circuit)
        clustered_circuit = serial_func(clustered_circuit)
        
    return clustered_circuit

def conduct_clustering_for_insert(circuit, max_cluster_degree):
    # serial_func = SerialMergeFunc(1)
    # clustered_circuit = serial_func(circuit)
    clustered_circuit = circuit
    
    for degree in range(2,max_cluster_degree+1):
        parallel_func=ParallelMergeFunc_full(degree)
        serial_func=SerialMergeFunc_moment(degree)
        
        clustered_circuit = parallel_func(clustered_circuit)
        clustered_circuit = serial_func(clustered_circuit)
        
    return clustered_circuit

def conduct_clustering_collapse(circuit, max_cluster_degree):
    serial_func = SerialMergeFunc_collapse(1)
    clustered_circuit = serial_func(circuit)
    # print(clustered_circuit)
    
    for degree in range(2,max_cluster_degree+1):
        parallel_func=ParallelMergeFunc_full(degree)
        serial_func=SerialMergeFunc_collapse(degree)
        
        clustered_circuit = parallel_func(clustered_circuit)
        # print(clustered_circuit)
        clustered_circuit = serial_func(clustered_circuit)
        # print(clustered_circuit)
        
    return clustered_circuit

def apply_BlinQ_module(circuit, max_cluster_degree):


    # if ENABLE_QSIM:    
    #     qsim_original = test_performance_original_qsim(circuit)
    #     return circuit

    max_cluster_degree = min(max_cluster_degree, len(circuit.all_qubits()))
    
    # clustered_circuit=conduct_clustering(circuit, max_cluster_degree)
    clustered_circuit=conduct_clustering_collapse(circuit, max_cluster_degree)

    ######### insert function ##############
    identity = True
    if identity:
        iter = max_cluster_degree
        
        # insert_func = InsertFrontAndBack_without_expand(max_cluster_degree)
        
        # insert_func = Insert_dual(max_cluster_degree)
        # inserted_circuit = insert_func(clustered_circuit)
        for _ in range(iter):
            insert_func = Insert_dual(max_cluster_degree)
            inserted_circuit = insert_func(clustered_circuit)
            clustered_circuit=conduct_clustering_for_insert(inserted_circuit, max_cluster_degree)
            
            # clustered_circuit=conduct_clustering(inserted_circuit, max_cluster_degree)
            # clustered_circuit=conduct_clustering_collapse(inserted_circuit, max_cluster_degree)

    
    return clustered_circuit

def get_sparsity(unitary, dimension):
    zero_count = 0
    for row in range(0, dimension):
        for col in range(0, dimension):
            if np.absolute(unitary[row][col]) == 0:
                zero_count += 1
                
    sparsity = zero_count / (dimension ** 2) * 100
    return sparsity
    
def get_size_gate_count(circuit, size):
    gate_count = 0
    for i, op in circuit.findall_operations(lambda op:len(op.qubits)==size):
        gate_count += 1
    return gate_count

def get_total_gate_count(circuit):
    gate_count = 0
    for i, op in circuit.findall_operations(lambda op:op):
        gate_count += 1
    return gate_count

def get_circuit_sparsity(circuit):
    sum=0
    for i, op in circuit.findall_operations(lambda op:op):
        op_sparsity = get_sparsity(cirq.unitary(op), len(op.qubits))
        sum += op_sparsity
    
    circuit_sparsity = sum / get_total_gate_count(circuit)    
    
    return circuit_sparsity

def get_size_sparsity(circuit, size):
    sum=0
    for i, op in circuit.findall_operations(lambda op:len(op.qubits)==size):
        op_sparsity = get_sparsity(cirq.unitary(op), len(op.qubits))
        sum += op_sparsity
    
    gate_sparsity = None
    if get_size_gate_count(circuit, size) != 0:
        gate_sparsity = sum / get_size_gate_count(circuit, size)
    
    return gate_sparsity

def get_circuit_depth(circuit):
    return len(circuit)

def get_gate_density(circuit):
    circuit_depth = get_circuit_depth(circuit)
    total_position = circuit_depth * len(circuit.all_qubits())
    total_pos_count = 0
    all_ops = circuit.findall_operations(lambda op:op)
    for i, op in all_ops:
        total_pos_count += len(op.qubits)
    gate_density = (total_pos_count / total_position) * 100
    
    return gate_density
    
def create_circuit_table(circuit, cluster_degree):
    
    result_circuit = apply_BlinQ_module(circuit, cluster_degree)
    
    original_total_gate = get_total_gate_count(circuit)
    reduced_total_gate = get_total_gate_count(result_circuit)
    
    original_circuit_depth = get_circuit_depth(circuit)
    reduced_circuit_depth = get_circuit_depth(result_circuit)
    
    original_gate_density = get_gate_density(circuit)
    reduced_gate_density = get_gate_density(result_circuit)
    
    original_sparsity = get_circuit_sparsity(circuit)
    reduced_sparsity = get_circuit_sparsity(result_circuit)
    
    return original_total_gate, reduced_total_gate, original_circuit_depth, reduced_circuit_depth, \
        original_gate_density, reduced_gate_density, original_sparsity, reduced_sparsity
    
def print_create_circuit_table(circuit, bench_name, num_qubits):
    gate_count_data_row = []
    circuit_depth_data_row = []
    gate_density_data_row = []
    circuit_sparsity_data_row = []

    for cd in range(2, 7, 1):
        original_total_gate, reduced_total_gate, original_circuit_depth, reduced_circuit_depth, \
            original_gate_density, reduced_gate_density, original_sparsity, reduced_sparsity \
                = create_circuit_table(circuit, cd)

        gate_count_algo_row = [bench_name, num_qubits, cd, original_total_gate, reduced_total_gate]
        circuit_depth_algo_row = [bench_name, num_qubits, cd, original_circuit_depth, reduced_circuit_depth]
        gate_density_algo_row = [bench_name, num_qubits, cd, original_gate_density, reduced_gate_density]
        circuit_sparsity_algo_row = [bench_name, num_qubits, cd, original_sparsity, reduced_sparsity]

        gate_count_data_row.append(gate_count_algo_row)
        circuit_depth_data_row.append(circuit_depth_algo_row)
        gate_density_data_row.append(gate_density_algo_row)
        circuit_sparsity_data_row.append(circuit_sparsity_algo_row)

    gate_count_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'original total gate', 'reduced total gate']
    circuit_depth_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'original circuit depth', 'reduced circuit depth']
    gate_density_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'original gate density', 'reduced gate density']
    circuit_sparsity_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'original sparsity', 'increased sparsity']

    gate_count_data_row.insert(0, gate_count_header_row)
    circuit_depth_data_row.insert(0, circuit_depth_header_row)
    gate_density_data_row.insert(0, gate_density_header_row)
    circuit_sparsity_data_row.insert(0, circuit_sparsity_header_row)

    print("gate count table")
    print(tabulate(gate_count_data_row))
    
    print("circuit depth table")
    print(tabulate(circuit_depth_data_row))
    
    print("gate density table")
    print(tabulate(gate_density_data_row))
    
    print("circuit sparsity table")
    print(tabulate(circuit_sparsity_data_row))
    
def create_gate_table(circuit, cluster_degree):
    result_circuit = apply_BlinQ_module(circuit, cluster_degree)
    
    original_gate_count_percenate_list=[0]*7
    optimized_gate_count_percenate_list=[0]*7
    original_size_gate_list = [0]*7
    optimized_size_gate_list = [0]*7
    original_size_sparsity_list = [0]*7
    optimized_size_sparsity_list = [0]*7
    
    for size in range(1, 7, 1):
        original_size_gate_list[size] = get_size_gate_count(circuit, size)
        optimized_size_gate_list[size] = get_size_gate_count(result_circuit, size)
        
        original_gate_count_percenate_list[size] = original_size_gate_list[size] / get_total_gate_count(circuit) * 100
        optimized_gate_count_percenate_list[size] = optimized_size_gate_list[size] / get_total_gate_count(result_circuit) * 100
    
        original_size_sparsity_list[size] = get_size_sparsity(circuit, size)
        optimized_size_sparsity_list[size] = get_size_sparsity(result_circuit, size)
        
    return original_gate_count_percenate_list, optimized_gate_count_percenate_list, original_size_gate_list, optimized_size_gate_list, \
        original_size_sparsity_list, optimized_size_sparsity_list
    
def print_create_gate_table(circuit, bench_name, num_qubits):
    size_count_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'target_size', 'original size gate count', 'changed size gate count']
    size_percentage_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'target_size', 'original size percentage', 'changed size percentage']
    size_sparsity_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'target_size', 'original size avg sparsity', 'changed size avg sparsity']

    for cd in range(2, 7, 1):
        original_gate_count_percenate_list, optimized_gate_count_percenate_list, \
            original_size_gate_list, optimized_size_gate_list, original_size_sparsity_list, optimized_size_sparsity_list\
                = create_gate_table(circuit, cd)

        size_count_data_row = []    
        size_percentage_data_row = []    
        size_sparsity_data_row = []    
        
        for size in range(1, 7, 1):
            size_count_algo_row = [bench_name, num_qubits, cd, size, original_size_gate_list[size], optimized_size_gate_list[size]]
            size_count_data_row.append(size_count_algo_row)
            
            size_percentage_algo_row = [bench_name, num_qubits, cd, size, original_gate_count_percenate_list[size], optimized_gate_count_percenate_list[size]]
            size_percentage_data_row.append(size_percentage_algo_row)
            
            size_sparsity_algo_row = [bench_name, num_qubits, cd, size, original_size_sparsity_list[size], optimized_size_sparsity_list[size]]
            size_sparsity_data_row.append(size_sparsity_algo_row)

        size_count_data_row.insert(0, size_count_header_row)
        size_percentage_data_row.insert(0, size_percentage_header_row)
        size_sparsity_data_row.insert(0, size_sparsity_header_row)

        print("size count table")
        print(tabulate(size_count_data_row))
        
        print("size percentage table")
        print(tabulate(size_percentage_data_row))
        
        print("size sparsity table")
        print(tabulate(size_sparsity_data_row))

def print_create_circuit_table_ppt(circuit, bench_name, num_qubits):
    gate_count_data_row = []
    circuit_depth_data_row = []
    gate_density_data_row = []
    circuit_sparsity_data_row = []

    for cd in range(2, 7, 1):
        original_total_gate, reduced_total_gate, original_circuit_depth, reduced_circuit_depth, \
            original_gate_density, reduced_gate_density, original_sparsity, reduced_sparsity \
                = create_circuit_table(circuit, cd)

        gate_count_algo_row = [bench_name, num_qubits, cd, original_total_gate, reduced_total_gate]
        circuit_depth_algo_row = [bench_name, num_qubits, cd, original_circuit_depth, reduced_circuit_depth]
        gate_density_algo_row = [bench_name, num_qubits, cd, original_gate_density, reduced_gate_density]
        circuit_sparsity_algo_row = [bench_name, num_qubits, cd, original_sparsity, reduced_sparsity]

        gate_count_data_row.append(gate_count_algo_row)
        circuit_depth_data_row.append(circuit_depth_algo_row)
        gate_density_data_row.append(gate_density_algo_row)
        circuit_sparsity_data_row.append(circuit_sparsity_algo_row)

    gate_count_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'original total gate', 'reduced total gate']
    circuit_depth_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'original circuit depth', 'reduced circuit depth']
    gate_density_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'original gate density', 'reduced gate density']
    circuit_sparsity_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'original sparsity', 'increased sparsity']

    gate_count_data_row.insert(0, gate_count_header_row)
    circuit_depth_data_row.insert(0, circuit_depth_header_row)
    gate_density_data_row.insert(0, gate_density_header_row)
    circuit_sparsity_data_row.insert(0, circuit_sparsity_header_row)


    prs = Presentation()
    sld=[]
    slide_num=4
    current_slide = 0
    for _ in range(slide_num):
        sld.append(prs.slides.add_slide(prs.slide_layouts[5]))
        
    sld[0].shapes[0].text = 'Gate Count Table'
    sld[1].shapes[0].text = 'Circuit Depth Table'
    sld[2].shapes[0].text = 'Gate Density Table'
    sld[3].shapes[0].text = 'Circuit Sparsity Table'

    rows = 6		# number of rows in table
    cols = 5		# number of columns in table
    
    # Add Shape object (table)
    # Arguments: number of rows, number of columns, x-coordinate of top-left corner, y-coordinate of top-left corner, width,  height
    table_shape=[]
    
    table_shape.append(sld[0].shapes.add_table(rows, cols, Cm(1), Cm(5), Cm(23), Cm(8)))
    table_shape.append(sld[1].shapes.add_table(rows, cols, Cm(1), Cm(5), Cm(23), Cm(8)))
    table_shape.append(sld[2].shapes.add_table(rows, cols, Cm(1), Cm(5), Cm(23), Cm(8)))
    table_shape.append(sld[3].shapes.add_table(rows, cols, Cm(1), Cm(5), Cm(23), Cm(8)))

    for data in [gate_count_data_row, circuit_depth_data_row, gate_density_data_row, circuit_sparsity_data_row]:
        
        table = table_shape[current_slide].table
        
        for index in range(len(data)):
            for index2 in range(len(data[0])):
                data[index][index2]=str(data[index][index2])
        table_row = data
        for row_index in range(len(table_row)):
            for col_index in range(len(table_row[0])):
                cell = table.cell(row_index, col_index)	 # Getting a cell object
                cell.text = table_row[row_index][col_index]   # Set the value with the text property
                pg = cell.text_frame.paragraphs[0]
                pg.font.color.rgb = RGBColor(0, 0, 0)  # Setting the font color for paragraphs
                # pg.font.size = Pt(15)		                # Setting the font size of paragraphs
                pg.aligment = PP_ALIGN.CENTER	
                
        current_slide += 1

    prs.save('circuit_summary_table.pptx')
    
def print_create_gate_table_ppt(circuit, bench_name, num_qubits):
    size_count_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'target_size', 'original size gate count', 'changed size gate count']
    size_percentage_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'target_size', 'original size percentage', 'changed size percentage']
    size_sparsity_header_row = ['algorithm', 'qubit_num', 'cluster_degree', 'target_size', 'original size avg sparsity', 'changed size avg sparsity']

    prs = Presentation()
    sld=[]
    current_slide=0
    slide_num=5 * 3
    for index in range(slide_num):
        sld.append(prs.slides.add_slide(prs.slide_layouts[5]))
        sld[index].shapes[0].text = "Table"
        
    rows = 7		# number of rows in table
    cols = 6		# number of columns in table
    
    table_shape=[]
    
    for index in range(slide_num):
        table_shape.append(sld[index].shapes.add_table(rows, cols, Cm(1), Cm(5), Cm(23), Cm(8)))
        

    for cd in range(2, 7, 1):
        original_gate_count_percenate_list, optimized_gate_count_percenate_list, \
            original_size_gate_list, optimized_size_gate_list, original_size_sparsity_list, optimized_size_sparsity_list\
                = create_gate_table(circuit, cd)

        size_count_data_row = []    
        size_percentage_data_row = []    
        size_sparsity_data_row = []    
        
        for size in range(1, 7, 1):
            size_count_algo_row = [bench_name, num_qubits, cd, size, original_size_gate_list[size], optimized_size_gate_list[size]]
            size_count_data_row.append(size_count_algo_row)
            
            size_percentage_algo_row = [bench_name, num_qubits, cd, size, original_gate_count_percenate_list[size], optimized_gate_count_percenate_list[size]]
            size_percentage_data_row.append(size_percentage_algo_row)
            
            size_sparsity_algo_row = [bench_name, num_qubits, cd, size, original_size_sparsity_list[size], optimized_size_sparsity_list[size]]
            size_sparsity_data_row.append(size_sparsity_algo_row)

        size_count_data_row.insert(0, size_count_header_row)
        size_percentage_data_row.insert(0, size_percentage_header_row)
        size_sparsity_data_row.insert(0, size_sparsity_header_row)


        for data in [size_count_data_row, size_percentage_data_row, size_sparsity_data_row]:
            table = table_shape[current_slide].table	# Create Table object        
            for index in range(len(data)):
                for index2 in range(len(data[0])):
                    data[index][index2]=str(data[index][index2])
                    
            table_row = data
            for row_index in range(len(table_row)):
                for col_index in range(len(table_row[0])):
                    cell = table.cell(row_index, col_index)	 # Getting a cell object
                    cell.text = table_row[row_index][col_index]   # Set the value with the text property
                    pg = cell.text_frame.paragraphs[0]
                    pg.font.color.rgb = RGBColor(0, 0, 0)  # Setting the font color for paragraphs
                    # pg.font.size = Pt(15)		                # Setting the font size of paragraphs
                    pg.aligment = PP_ALIGN.CENTER	
                    
            current_slide += 1
                    
    prs.save('gate_summary_table.pptx')



bench_circuit = circuit_from_qasm(qasm_input)

# print_create_circuit_table(bench_circuit, bench_name, num_qubits)

# print_create_gate_table(bench_circuit, bench_name, num_qubits)

print_create_circuit_table_ppt(bench_circuit, bench_name, num_qubits)

print_create_gate_table_ppt(bench_circuit, bench_name, num_qubits)







